import io
import datetime
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from utils import get_gemini_response

def clean_ai_text(text: str) -> str:
    """
    Removes weird ASCII control characters, all asterisks (*), and backticks (`).
    Preserves normal punctuation, newlines, and other symbols.

    Steps:
      1) Strip ASCII control chars [\x00-\x08, \x0B-\x0C, \x0E-\x1F, \x7F-\x9F].
      2) Remove carriage returns (\r).
      3) Collapse multiple blank lines into one.
      4) Remove all asterisks (*) and backticks (`).
      5) Trim leading/trailing whitespace.
    """
    # 1. Remove control chars
    cleaned = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F-\x9F]', '', text)
    
    # 2. Remove carriage returns
    cleaned = cleaned.replace('\r', '')
    
    # 3. Collapse multiple blank lines
    cleaned = re.sub(r'\n\s*\n+', '\n\n', cleaned)
    
    # 4. Remove asterisks and backticks
    cleaned = re.sub(r'[`*]', '', cleaned)
    
    # 5. Trim whitespace
    cleaned = cleaned.strip()
    
    return cleaned

def generate_eda_report_ppt(
    eda_metadata,
    df,
    numeric_figs=None,
    categorical_figs=None,
    correlation_figs=None,
    time_series_figs=None,
    outlier_figs=None,
    dataset_name="Dataset.csv"
):
    """
    Generates a PPTX report with a dark background and white text,
    splitting commentary by both line count (max ~22 lines) and word count (max ~300 words).
    If either limit is exceeded, we start a new slide.

    1) Title slide (dark background).
    2) Overview slide (rows, columns).
    3) Per-section commentary + figure slides.
    4) Graceful kaleido error handling for figures.
    5) Conclusion.

    Returns a BytesIO with the PPTX content.
    """

    # Convert None to empty lists
    numeric_figs = numeric_figs or []
    categorical_figs = categorical_figs or []
    correlation_figs = correlation_figs or []
    time_series_figs = time_series_figs or []
    outlier_figs = outlier_figs or []

    numeric_insights = get_gemini_response(f""" Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Summarize the **numeric columns** in a friendly, user-focused way.
                                            - Write about 150 to 250 words total.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Organize the summary in **bullet points** (1–2 lines each).
                                            - For each bullet, give a short explanation suitable for **non-technical** users.
                                            - Include each numeric column’s typical range, average, or any key outliers or patterns.
                                            - If the dataset has no numeric columns, say “No numeric columns found.”
                                            - Output must be the **final text only** no formatting.
                                            """)
    categorical_insights = get_gemini_response(f""" Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Summarize the **categorical columns** in a friendly, user-focused way.
                                            - Write about 150 to 250 words total.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Organize the summary in **bullet points** (1–2 lines each).
                                            - For each bullet, give a short explanation suitable for **non-technical** users.
                                            - Include each numeric column’s typical range, average, or any key outliers or patterns.
                                            - If the dataset has no categorical columns, say “No categorical columns found.”
                                            - Output must be the **final text only** no formatting.
                                            """)
    correlation_insights = get_gemini_response(f""" Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Summarize the **correlation columns** in a friendly, user-focused way.
                                            - Write about 150 to 250 words total.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Organize the summary in **bullet points** (1–2 lines each).
                                            - For each bullet, give a short explanation suitable for **non-technical** users.
                                            - Include each numeric column’s typical range, average, or any key outliers or patterns.
                                            - If the dataset has no correlation columns, say “No correlation columns found.”
                                            - Output must be the **final text only** no formatting.
                                            """)
    outlier_insights = get_gemini_response(f""" Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Summarize the **outliners columns** in a friendly, user-focused way.
                                            - Write about 150 to 250 words total.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Organize the summary in **bullet points** (1–2 lines each).
                                            - For each bullet, give a explanation suitable for **non-technical** users.
                                            - Include each numeric column’s typical range, average, or any key outliers or patterns.
                                            - If the dataset has no outliners columns, say “No outliners columns found.”
                                            - Output must be the **final text only** no formatting.
                                            """)
    time_series_insights = get_gemini_response(f""" Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Summarize the **time series columns** in a friendly, user-focused way.
                                            - Write about 150 to 250 words total.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Organize the summary in **bullet points** (1–2 lines each).
                                            - For each bullet, give a short explanation suitable for **non-technical** users.
                                            - Include each numeric column’s typical range, average, or any key outliers or patterns.
                                            - If the dataset has no time series columns, say “No time series columns found.”
                                            - Output must be the **final text only** no formatting.
                                            """)
    overall_insights = get_gemini_response(f"""Here is the dataset context (in JSON):
                                            {eda_metadata}

                                            INSTRUCTIONS:
                                            - Provide an **depth summary** of the entire EDA in a friendly, user-focused way and breif enough.
                                            - Use plain text only—no code blocks, no triple backticks, no excessive Markdown.
                                            - Focus on **actionable insights**, key findings, or interesting patterns across numeric, categorical, correlation, or time-series data.
                                            - Avoid repeating trivial details; highlight the big takeaways that **non-technical** readers can understand.
                                            - Output must be the **final text only**, no formatting.
                                            """)
    
    # Clean AI-generated text
    numeric_insights = clean_ai_text(numeric_insights)
    categorical_insights = clean_ai_text(categorical_insights)
    correlation_insights = clean_ai_text(correlation_insights)
    outlier_insights = clean_ai_text(outlier_insights)
    time_series_insights = clean_ai_text(time_series_insights)
    overall_insights = clean_ai_text(overall_insights)



    # Create the presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    #### 1) Dark Slide Helper ####
    def add_dark_slide(prs):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        shape = slide.shapes.add_shape(
            autoshape_type_id=1,  # rectangle
            left=Inches(0), top=Inches(0),
            width=prs.slide_width, height=prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(26, 42, 58)  # #1A2A3A
        shape.line.fill.background()
        return slide

    #### 2) Text Box Helper ####
    def add_textbox(slide, left, top, width, height,
                    text, font_size=24,
                    color=RGBColor(250,250,250),
                    align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        return box

    #### 3) Hybrid Chunking: By Lines & Word Count ####
    def chunk_text_hybrid(text, max_lines=15, max_words=140):
        """
        Splits text into chunks, each chunk having at most `max_lines` lines
        and `max_words` total words. We treat each newline as a distinct line.

        Steps:
          1) Split the text by actual newlines -> lines
          2) For each line, count words
          3) Keep appending lines until we exceed max_lines or max_words
          4) Start a new chunk

        Returns a list of multiline strings (each chunk = one slide).
        """

        raw_lines = text.split('\n')
        chunks = []
        chunk_lines = []
        chunk_word_count = 0

        for line in raw_lines:
            words_in_line = len(line.split())
            # If adding this line would exceed lines or words, start a new chunk
            if (len(chunk_lines) + 1 > max_lines) or (chunk_word_count + words_in_line > max_words):
                # push current chunk
                chunks.append("\n".join(chunk_lines))
                chunk_lines = [line]
                chunk_word_count = words_in_line
            else:
                chunk_lines.append(line)
                chunk_word_count += words_in_line

        # leftover
        if chunk_lines:
            chunks.append("\n".join(chunk_lines))

        return chunks

    #### 4) Add Section ####
    def add_section(section_title, commentary_text, figs):
        # chunk commentary
        commentary_chunks = chunk_text_hybrid(commentary_text, max_lines=15, max_words=140)

        # Create bullet slides for commentary
        for idx, chunk in enumerate(commentary_chunks):
            slide = add_dark_slide(prs)
            add_textbox(
                slide,
                Inches(0.5), Inches(0.3),
                Inches(12), Inches(0.5),
                text=section_title if idx == 0 else f"{section_title} (cont.)",
                font_size=28
            )
            add_textbox(
                slide,
                Inches(0.5), Inches(1.2),
                Inches(12), Inches(5),
                text=chunk,
                font_size=18
            )

        # For each figure, add a new slide
        for i, fig in enumerate(figs):
            slide = add_dark_slide(prs)
            add_textbox(
                slide,
                Inches(0.5), Inches(0.3),
                Inches(12), Inches(0.5),
                text=f"{section_title} (Fig {i+1})",
                font_size=24
            )
            try:
                img_bytes = fig.to_image(format="png")
                left = (prs.slide_width - Inches(8)) / 2
                top = Inches(1.3)
                slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=Inches(8))
            except Exception as e:
                err = str(e).lower()
                if "kaleido" in err:
                    msg = ("Could not export figure because kaleido is not installed.\n"
                           "Install with: pip install -U kaleido")
                else:
                    msg = f"Error rendering figure {i+1}: {e}"
                add_textbox(
                    slide,
                    Inches(1), Inches(1.5),
                    Inches(10), Inches(2),
                    text=msg,
                    font_size=16
                )

    #### 5) Build Slides ####

    # Title Slide
    slide = add_dark_slide(prs)
    add_textbox(
        slide,
        Inches(0.5), Inches(0.5),
        Inches(12), Inches(1.0),
        text="EDA Pro Analysis Report",
        font_size=36, align=PP_ALIGN.CENTER
    )
    date_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
    add_textbox(
        slide,
        Inches(0.5), Inches(1.5),
        Inches(12), Inches(1.0),
        text=f"Dataset: {dataset_name}\nGenerated on {date_str}",
        font_size=18, align=PP_ALIGN.CENTER
    )

    # Overview Slide
    slide = add_dark_slide(prs)
    add_textbox(
        slide,
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.5),
        text="Dataset Overview",
        font_size=28
    )
    row_count = df.shape[0]
    col_count = df.shape[1]
    lines = [f"Rows: {row_count}", f"Columns: {col_count}", ""]
    if eda_metadata and "columns" in eda_metadata:
        lines.append("Columns:")
        for col_name in eda_metadata["columns"].keys():
            lines.append(f"• {col_name}")
    overview_text = "\n".join(lines)
    add_textbox(
        slide,
        Inches(0.5), Inches(1.2),
        Inches(12), Inches(5),
        text=overview_text,
        font_size=18
    )

    # Numeric
    add_section("Numeric Analysis", numeric_insights, numeric_figs)
    # Categorical
    add_section("Categorical Analysis", categorical_insights, categorical_figs)
    # Correlation
    add_section("Correlation Analysis", correlation_insights, correlation_figs)
    # Time Series
    add_section("Time Series Analysis", time_series_insights, time_series_figs)
    # Outliers
    add_section("Outlier Analysis", outlier_insights, outlier_figs)

    # Overall Insights
    from_texts = chunk_text_hybrid(overall_insights, max_lines=18, max_words=140)
    for idx, chunk in enumerate(from_texts):
        slide = add_dark_slide(prs)
        add_textbox(
            slide,
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.5),
            text="Overall Data Whisperer Insights" if idx == 0 else "Overall Data Whisperer Insights (cont.)",
            font_size=28
        )
        add_textbox(
            slide,
            Inches(0.5), Inches(1.2),
            Inches(12), Inches(5),
            text=chunk,
            font_size=18
        )

    # Conclusion Slide
    slide = add_dark_slide(prs)
    add_textbox(
        slide,
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.5),
        text="Conclusion",
        font_size=28
    )
    add_textbox(
        slide,
        Inches(0.5), Inches(1.2),
        Inches(12), Inches(5),
        text=(
            "• This concludes the Data Whisperer analysis report.\n"
            "• We hope these insights help guide your next steps!"
        ),
        font_size=18
    )

    # Save to BytesIO
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer
